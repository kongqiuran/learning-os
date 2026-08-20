# 文件说明：把上传文件抽成纯文本，供 DocumentAnalysis 的 LLM 分析使用。
# Path 来自 Python 标准库 pathlib；fitz 来自 PyMuPDF，用来读取 PDF；Presentation 来自 python-pptx，用来读取 PPTX。
from pathlib import Path

import fitz
from pptx import Presentation

from src.config import BASE_DIR


SUPPORTED_MIME_TYPES = {
    "application/pdf": "PDF",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "PPTX",
    "text/plain": "TXT",
    "text/markdown": "MD",
    "text/x-markdown": "MD",
}


def extract_text(file_path, mime_type):
    # extract_text 是非视觉流水线的文本抽取入口。
    # PDF/PPTX/TXT/MD 最终都要变成字符串，后续 analyze_document 才能把文本发给 LLM。
    path = Path(file_path)
    if not path.is_absolute():
        path = BASE_DIR / path

    source_type = get_source_type(path, mime_type)
    if source_type == "PDF":
        text = _extract_pdf(path)
    elif source_type == "PPTX":
        text = _extract_pptx(path)
    else:
        text = path.read_text(encoding="utf-8-sig")

    if not text.strip():
        raise ValueError(f"No text was found in the {source_type} file. OCR is not supported.")
    return text


def get_source_type(file_path, mime_type):
    # source_type 是给 AI 分析器看的来源类型，例如 PDF、PPTX、TXT、MD。
    # MD/TXT 优先看扩展名；PDF/PPTX 用 MIME 类型判断，和上传校验保持一致。
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix == ".md":
        return "MD"
    if suffix == ".txt":
        return "TXT"
    source_type = SUPPORTED_MIME_TYPES.get((mime_type or "").lower().strip())
    if source_type is None:
        raise ValueError("Supported file types are PDF, PPTX, TXT, and MD.")
    return source_type


def _extract_pdf(path):
    # fitz.open 来自 PyMuPDF，它会打开 PDF；page.get_text("text") 表示抽取页面中的文本层。
    # 如果 PDF 是扫描图片且没有文本层，这里拿不到 OCR 内容，后续会报“没有文本”。
    with fitz.open(path) as document:
        pages = [page.get_text("text").strip() for page in document]
    return "\n\n".join(page for page in pages if page)


def _extract_pptx(path):
    # Presentation 来自 python-pptx；presentation.slides 是每一页幻灯片。
    # shape.text 是 PPT 形状里的文字框内容，hasattr 用来确认这个形状确实有 text 属性。
    presentation = Presentation(path)
    slides = []
    for slide in presentation.slides:
        parts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if parts:
            slides.append("\n".join(parts))
    return "\n\n".join(slides)
