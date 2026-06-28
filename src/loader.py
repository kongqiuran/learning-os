from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_files(files):
    documents = []
    for file in files:
        text = extract_text_from_file(file)
        if text.strip():
            documents.append(
                {
                    "filename": getattr(file, "name", "unknown"),
                    "text": text.strip(),
                }
            )
    return documents


def extract_text_from_file(file):
    filename = getattr(file, "name", str(file))
    suffix = Path(filename).suffix.lower()
    data = _read_file_bytes(file)

    if suffix == ".pdf":
        return extract_pdf_text(data)
    if suffix == ".pptx":
        return extract_pptx_text(data)
    if suffix in {".txt", ".md"}:
        return data.decode("utf-8", errors="ignore")

    return ""


def extract_pdf_text(data):
    reader = PdfReader(BytesIO(data))
    pages = []

    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"## Page {index}\n{text}")

    return "\n\n".join(pages)


def extract_pptx_text(data):
    presentation = Presentation(BytesIO(data))
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())

        if lines:
            slides.append(f"## Slide {index}\n" + "\n".join(lines))

    return "\n\n".join(slides)


def _read_file_bytes(file):
    if hasattr(file, "read"):
        file.seek(0)
        return file.read()
    return Path(file).read_bytes()
