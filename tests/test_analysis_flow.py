import tempfile
import unittest
from pathlib import Path

import fitz
from pptx import Presentation


TEST_DIR = tempfile.TemporaryDirectory()

from src.database import create_database_tables, get_db_session  # noqa: E402
from src.ai.llm_client import LLMGenerationError  # noqa: E402
from src.models import Chapter, Course, Document, DocumentAnalysis, LearningPackage, User  # noqa: E402
from src.services.analysis_service import (  # noqa: E402
    _validate_document_course_binding,
    _load_course_documents,
    analyze_course,
    create_learning_package_task,
    get_learning_package,
    get_scoped_packages,
)
from src.services.file_parser_service import extract_text  # noqa: E402
from src.services.document_service import (  # noqa: E402
    delete_document_for_user,
    list_documents_for_course,
    save_uploaded_document,
)


class FakeUploadedFile:
    def __init__(self, name, mime_type, data):
        self.name = name
        self.type = mime_type
        self._data = data

    def getvalue(self):
        return self._data


class FakeLLMClient:
    def generate(self, system_prompt, user_prompt, stage="unknown"):
        if '"question_patterns"' in system_prompt:
            return {
                "summary": "Test summary",
                "topics": [
                    {
                        "name": "Core topic",
                        "importance": 5,
                        "exam_value": "Common calculation topic",
                        "reason": "Repeated in the material",
                    }
                ],
                "formulas": [],
                "question_patterns": [],
                "errors": [],
            }
        if '"formula_inventory"' in system_prompt:
            return {
                "knowledge_map": {"Core topic": []},
                "chapter_relations": [],
                "priority_ranking": ["Core topic"],
                "formula_inventory": [],
                "exam_focus": [],
            }
        return {
            "course_map": {"Core topic": []},
            "chapter_summary": ["Chapter summary"],
            "key_points": ["Core topic"],
            "formula_book": [],
            "exam_focus": [
                {
                    "topic": "Core topic",
                    "importance": 5,
                    "reason": "Central course concept",
                    "evidence": ["来源于课程资料分析"],
                }
            ],
            "questions": [{"question": "Test?", "answer": "Yes"}],
            "exam_strategy": {
                "priority_order": ["Core topic"],
                "study_advice": "Review the highest priority topic first.",
            },
        }


class StageRecordingLLMClient(FakeLLMClient):
    def __init__(self):
        self.stages = []

    def generate(self, system_prompt, user_prompt, stage="unknown"):
        self.stages.append(stage)
        return super().generate(system_prompt, user_prompt, stage)


class FailingLLMClient:
    def generate(self, system_prompt, user_prompt, stage="unknown"):
        raise LLMGenerationError(
            "Invalid structured output.",
            stage=stage,
            retry_count=2,
            error_type="JSONParseError",
            response_preview="broken output",
        )


class AnalysisFlowTest(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        create_database_tables()

    @classmethod
    def tearDownClass(cls):
        TEST_DIR.cleanup()

    def setUp(self):
        with get_db_session() as session:
            for model in (LearningPackage, DocumentAnalysis, Document, Chapter, Course, User):
                session.query(model).delete()

    def test_follow_documents_are_scoped_to_one_chapter(self):
        first_path = Path(TEST_DIR.name) / "chapter-1.txt"
        second_path = Path(TEST_DIR.name) / "chapter-2.txt"
        first_path.write_text("First chapter", encoding="utf-8")
        second_path.write_text("Second chapter", encoding="utf-8")
        with get_db_session() as session:
            user = User(email="chapters@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Scoped course")
            session.add(course)
            session.flush()
            first_chapter = Chapter(course_id=course.id, title="第一章", position=0)
            second_chapter = Chapter(course_id=course.id, title="第二章", position=1)
            session.add_all([first_chapter, second_chapter])
            session.flush()
            first_document = self._document(user.id, course.id, first_path, "text/plain")
            first_document.document_type = "NOTES"
            first_document.chapter_id = first_chapter.id
            second_document = self._document(user.id, course.id, second_path, "text/plain")
            second_document.document_type = "NOTES"
            second_document.chapter_id = second_chapter.id
            session.add_all([first_document, second_document])
            session.flush()
            user_id, course_id, first_chapter_id = user.id, course.id, first_chapter.id

        _, documents = _load_course_documents(course_id, user_id, "follow", scope_chapter_id=first_chapter_id)

        self.assertEqual([document.original_filename for document in documents], ["chapter-1.txt"])

    def test_follow_chapter_generation_skips_course_wide_pass(self):
        path = Path(TEST_DIR.name) / "focused-chapter.txt"
        path.write_text("Focused chapter material", encoding="utf-8")
        with get_db_session() as session:
            user = User(email="fast-follow@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Fast follow")
            session.add(course)
            session.flush()
            chapter = Chapter(course_id=course.id, title="第一章", position=0)
            session.add(chapter)
            session.flush()
            document = self._document(user.id, course.id, path, "text/plain")
            document.document_type = "NOTES"
            document.chapter_id = chapter.id
            session.add(document)
            session.flush()
            user_id, course_id, chapter_id = user.id, course.id, chapter.id

        task = create_learning_package_task(course_id, user_id, "follow", scope_chapter_id=chapter_id)
        client = StageRecordingLLMClient()
        result = analyze_course(course_id, user_id, llm_client=client, package_id=task.id, scene="follow", scope_chapter_id=chapter_id)

        self.assertEqual(result.status, "completed")
        self.assertIn("follow_chapter_generator", client.stages)
        self.assertNotIn("course_analyzer", client.stages)
        self.assertNotIn("learning_package_generator", client.stages)

    def test_chapter_packages_have_independent_versions_and_stale_detection(self):
        first_path = Path(TEST_DIR.name) / "scope-version-1.txt"
        second_path = Path(TEST_DIR.name) / "scope-version-2.txt"
        first_path.write_text("First scope", encoding="utf-8")
        second_path.write_text("Second scope", encoding="utf-8")
        with get_db_session() as session:
            user = User(email="scope-versions@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Scope versions")
            session.add(course)
            session.flush()
            first_chapter = Chapter(course_id=course.id, title="First", position=0)
            second_chapter = Chapter(course_id=course.id, title="Second", position=1)
            session.add_all([first_chapter, second_chapter])
            session.flush()
            first_document = self._document(user.id, course.id, first_path, "text/plain")
            first_document.document_type = "NOTES"
            first_document.chapter_id = first_chapter.id
            second_document = self._document(user.id, course.id, second_path, "text/plain")
            second_document.document_type = "NOTES"
            second_document.chapter_id = second_chapter.id
            session.add_all([first_document, second_document])
            session.flush()
            user_id = user.id
            course_id = course.id
            first_chapter_id = first_chapter.id
            second_chapter_id = second_chapter.id
            first_document_id = first_document.id

        first = create_learning_package_task(course_id, user_id, "follow", scope_chapter_id=first_chapter_id)
        second = create_learning_package_task(course_id, user_id, "follow", scope_chapter_id=second_chapter_id)
        first_again = create_learning_package_task(course_id, user_id, "follow", scope_chapter_id=first_chapter_id)

        self.assertEqual((first.version, second.version, first_again.version), (1, 1, 2))
        self.assertEqual(first.scope_key, f"chapter:{first_chapter_id}")
        chapter_packages, _ = get_scoped_packages(course_id, user_id)
        self.assertFalse(chapter_packages[str(first_chapter_id)].is_stale)

        with get_db_session() as session:
            session.get(Document, first_document_id).chapter_id = None

        chapter_packages, _ = get_scoped_packages(course_id, user_id)
        self.assertTrue(chapter_packages[str(first_chapter_id)].is_stale)

    def test_course_analysis_persists_results(self):
        pdf_path = Path(TEST_DIR.name) / "course.pdf"
        with fitz.open() as pdf:
            page = pdf.new_page()
            page.insert_text((72, 72), "Core topic and formula F = ma")
            pdf.save(pdf_path)

        with get_db_session() as session:
            user = User(email="student@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Physics")
            session.add(course)
            session.flush()
            document = Document(
                user_id=user.id,
                course_id=course.id,
                original_filename="course.pdf",
                stored_filename="course.pdf",
                file_path=str(pdf_path),
                mime_type="application/pdf",
                file_size=pdf_path.stat().st_size,
                document_type="TEXTBOOK",
            )
            session.add(document)
            session.flush()
            user_id, course_id = user.id, course.id

        package = analyze_course(course_id, user_id, llm_client=FakeLLMClient())
        self.assertEqual(package.status, "completed")
        self.assertIn("course_map", package.content_json)
        latest = get_learning_package(course_id, user_id)
        self.assertEqual(latest.id, package.id)

        with get_db_session() as session:
            self.assertEqual(session.query(DocumentAnalysis).count(), 1)
            self.assertEqual(session.query(LearningPackage).count(), 1)

    def test_two_courses_are_analyzed_in_isolation(self):
        signal_pdf = Path(TEST_DIR.name) / "CH1.pdf"
        verilog_pdf = Path(TEST_DIR.name) / "verilog.pdf"
        self._make_pdf(signal_pdf, "Signal and Systems convolution")
        self._make_pdf(verilog_pdf, "Verilog finite state machine")

        with get_db_session() as session:
            user = User(email="isolation@example.com", password_hash="test")
            session.add(user)
            session.flush()
            signal_course = Course(user_id=user.id, name="信号与系统")
            verilog_course = Course(user_id=user.id, name="Verilog")
            session.add_all([signal_course, verilog_course])
            session.flush()
            session.add_all(
                [
                    self._document(user.id, signal_course.id, signal_pdf),
                    self._document(user.id, verilog_course.id, verilog_pdf),
                ]
            )
            session.flush()
            user_id = user.id
            signal_course_id = signal_course.id
            verilog_course_id = verilog_course.id

        analyze_course(signal_course_id, user_id, llm_client=FakeLLMClient())
        with get_db_session() as session:
            analyzed_document_ids = {
                item.document_id for item in session.query(DocumentAnalysis).all()
            }
            signal_document_id = session.query(Document.id).filter_by(course_id=signal_course_id).scalar()
            verilog_document_id = session.query(Document.id).filter_by(course_id=verilog_course_id).scalar()
        self.assertIn(signal_document_id, analyzed_document_ids)
        self.assertNotIn(verilog_document_id, analyzed_document_ids)

    def test_pending_task_is_reused_and_completed(self):
        pdf_path = Path(TEST_DIR.name) / "queued.pdf"
        self._make_pdf(pdf_path, "Queued generation content")
        with get_db_session() as session:
            user = User(email="queued@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Queued course")
            session.add(course)
            session.flush()
            session.add(self._document(user.id, course.id, pdf_path))
            session.flush()
            user_id, course_id = user.id, course.id

        task = create_learning_package_task(course_id, user_id)
        self.assertEqual(task.status, "pending")
        result = analyze_course(
            course_id,
            user_id,
            llm_client=FakeLLMClient(),
            package_id=task.id,
        )

        self.assertEqual(result.id, task.id)
        self.assertEqual(result.status, "completed")
        with get_db_session() as session:
            self.assertEqual(session.query(LearningPackage).count(), 1)

    def test_failed_task_persists_stage_retry_and_error_detail(self):
        pdf_path = Path(TEST_DIR.name) / "failure.pdf"
        self._make_pdf(pdf_path, "Failure diagnostics content")
        with get_db_session() as session:
            user = User(email="failure@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Failure course")
            session.add(course)
            session.flush()
            session.add(self._document(user.id, course.id, pdf_path))
            session.flush()
            user_id, course_id = user.id, course.id

        with self.assertRaises(LLMGenerationError):
            analyze_course(course_id, user_id, llm_client=FailingLLMClient())

        package = get_learning_package(course_id, user_id)
        self.assertEqual(package.status, "failed")
        self.assertEqual(package.current_stage, "document_analyzer")
        self.assertEqual(package.retry_count, 2)
        self.assertEqual(package.error_type, "JSONParseError")
        self.assertIn("broken output", package.error_detail)

    def test_pdf_and_pptx_enter_analysis_flow(self):
        pdf_path = Path(TEST_DIR.name) / "formats.pdf"
        pptx_path = Path(TEST_DIR.name) / "formats.pptx"
        self._make_pdf(pdf_path, "PDF course material")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "PPTX course material"
        presentation.save(pptx_path)

        with get_db_session() as session:
            user = User(email="formats@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="Mixed formats")
            session.add(course)
            session.flush()
            session.add(self._document(user.id, course.id, pdf_path, "application/pdf"))
            session.add(
                self._document(
                    user.id,
                    course.id,
                    pptx_path,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            )
            session.flush()
            user_id, course_id = user.id, course.id

        package = analyze_course(course_id, user_id, llm_client=FakeLLMClient())
        self.assertEqual(package.status, "completed")
        with get_db_session() as session:
            self.assertEqual(
                session.query(DocumentAnalysis)
                .join(Document)
                .filter(Document.course_id == course_id)
                .count(),
                2,
            )
        self.assertIn("PPTX course material", extract_text(pptx_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation"))

    def test_document_course_mismatch_is_rejected(self):
        mismatched = type("DocumentStub", (), {"course_id": 2})()
        with self.assertRaisesRegex(ValueError, "Document-course mismatch detected\\."):
            _validate_document_course_binding(1, [mismatched])

    def test_txt_and_markdown_are_parsed_as_plain_text(self):
        txt_path = Path(TEST_DIR.name) / "notes.txt"
        md_path = Path(TEST_DIR.name) / "notes.md"
        txt_path.write_text("Plain course notes", encoding="utf-8")
        md_path.write_text("# Markdown course notes", encoding="utf-8")

        self.assertEqual(extract_text(txt_path, "text/plain"), "Plain course notes")
        self.assertEqual(extract_text(md_path, "text/markdown"), "# Markdown course notes")

    def test_material_inbox_types_multi_upload_and_delete(self):
        pdf_path = Path(TEST_DIR.name) / "inbox.pdf"
        pptx_path = Path(TEST_DIR.name) / "inbox.pptx"
        self._make_pdf(pdf_path, "Signals and Systems chapter one")
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = "Chapter 01"
        presentation.save(pptx_path)

        with get_db_session() as session:
            user = User(email="inbox@example.com", password_hash="test")
            session.add(user)
            session.flush()
            course = Course(user_id=user.id, name="信号与系统")
            session.add(course)
            session.flush()
            user_id, course_id = user.id, course.id

        uploads = [
            ("TEXTBOOK", FakeUploadedFile("CH1.pdf", "application/pdf", pdf_path.read_bytes())),
            (
                "SLIDES",
                FakeUploadedFile(
                    "Chapter01.pptx",
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    pptx_path.read_bytes(),
                ),
            ),
            ("EXAM", FakeUploadedFile("A卷.pdf", "application/pdf", pdf_path.read_bytes())),
            ("NOTES", FakeUploadedFile("note-1.txt", "text/plain", b"first note")),
            ("NOTES", FakeUploadedFile("note-2.txt", "text/plain", b"second note")),
        ]
        for document_type, uploaded_file in uploads:
            save_uploaded_document(
                user_id,
                course_id,
                uploaded_file,
                document_type=document_type,
            )

        documents = list_documents_for_course(user_id, course_id)
        types_by_name = {document.original_filename: document.document_type for document in documents}
        self.assertEqual(types_by_name["CH1.pdf"], "TEXTBOOK")
        self.assertEqual(types_by_name["Chapter01.pptx"], "SLIDES")
        self.assertEqual(types_by_name["A卷.pdf"], "EXAM")
        self.assertEqual(types_by_name["note-1.txt"], "NOTES")
        self.assertEqual(types_by_name["note-2.txt"], "NOTES")

        deleted_document = next(item for item in documents if item.original_filename == "note-1.txt")
        self.assertTrue(delete_document_for_user(deleted_document.id, user_id, course_id))
        remaining_names = {
            item.original_filename for item in list_documents_for_course(user_id, course_id)
        }
        self.assertNotIn("note-1.txt", remaining_names)
        for document in list_documents_for_course(user_id, course_id):
            delete_document_for_user(document.id, user_id, course_id)

    @staticmethod
    def _make_pdf(path, text):
        with fitz.open() as pdf:
            page = pdf.new_page()
            page.insert_text((72, 72), text)
            pdf.save(path)

    @staticmethod
    def _document(user_id, course_id, path, mime_type="application/pdf"):
        return Document(
            user_id=user_id,
            course_id=course_id,
            original_filename=path.name,
            stored_filename=path.name,
            file_path=str(path),
            mime_type=mime_type,
            file_size=path.stat().st_size,
            document_type="TEXTBOOK",
        )


if __name__ == "__main__":
    unittest.main()
